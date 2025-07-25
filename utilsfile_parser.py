import fitz  # PyMuPDF
import docx
import pandas as pd
import pptx
import os

def parse_pdf(file_path):
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text

def parse_docx(file_path):
    doc = docx.Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def parse_csv(file_path):
    df = pd.read_csv(file_path)
    return df.to_string(index=False)

def parse_pptx(file_path):
    prs = pptx.Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def parse_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()

def parse_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return parse_pdf(file_path)
    elif ext == ".docx":
        return parse_docx(file_path)
    elif ext == ".csv":
        return parse_csv(file_path)
    elif ext == ".pptx":
        return parse_pptx(file_path)
    elif ext in [".txt", ".md"]:
        return parse_txt(file_path)
    else:
        return "Unsupported file type"
